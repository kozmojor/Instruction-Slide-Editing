#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
preprocess.py

Build a *mismatched* & *random-layout* PPTX for an instructed-GPT layout-fixing demo.
- Collect N images and N captions *independently* from COCO (from HuggingFace: yerevann/coco-karpathy).
- Randomly pair them (mismatch by design).
- Place image & caption at random positions/sizes (likely overlap/margin breaches).

For demo purpose, this script only output one .pptx file following logic above. This file will be used in the instruction fix jupyter notebook. (See in slidedit.ipynb)
"""

import argparse
import io
import os
import random
from pathlib import Path
from typing import List, Dict, Any, Tuple

import requests
from datasets import load_dataset
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

EMU_PER_IN, DPI = 914400, 96
def emu2px(v): return int(round(float(v)/EMU_PER_IN*DPI))
def px2emu(v): return int(round(float(v)/DPI*EMU_PER_IN))

def parse_scale_range(s: str) -> Tuple[float, float]:
    try:
        a, b = s.split(",")
        a, b = float(a), float(b)
        if not (0 < a < 1 and 0 < b < 1):
            raise ValueError
        return (min(a, b), max(a, b))
    except Exception:
        raise argparse.ArgumentTypeError("scale-range must be like '0.22,0.60' within (0,1).")

def parse_font_candidates(s: str) -> Tuple[int, ...]:
    try:
        vals = tuple(int(x.strip()) for x in s.split(",") if x.strip())
        assert all(v >= 6 for v in vals)
        return vals
    except Exception:
        raise argparse.ArgumentTypeError("font-candidates must be like '26,32,36,40' (pt).")

def _rand_size(SW: int, SH: int, scale_range: Tuple[float, float]) -> Tuple[int, int]:
    wf = random.uniform(*scale_range)
    hf = random.uniform(*scale_range)
    return max(80, int(SW*wf)), max(80, int(SH*hf))

def _rand_pos(SW: int, SH: int, w: int, h: int) -> Tuple[int, int]:
    return (random.randint(0, max(1, SW - w)), random.randint(0, max(1, SH - h)))

def _force_overlap(ax, ay, aw, ah, bw, bh, SW, SH):
    """Return (x, y) for B to overlap A."""
    ox = random.randint(max(0, ax - bw//2), min(SW - bw, ax + aw - bw//2))
    oy = random.randint(max(0, ay - bh//2), min(SH - bh, ay + ah - bh//2))
    return ox, oy

def _force_margin_breach(w, h, SW, SH, side: str):
    """Return (x,y) that breaches margins or goes slightly OOB."""
    if side == "left":
        return -int(w*0.1), random.randint(0, SH - h)
    if side == "top":
        return random.randint(0, SW - w), -int(h*0.1)
    if side == "right":
        return SW - int(w*0.8), random.randint(0, SH - h)
    if side == "bottom":
        return random.randint(0, SW - w), SH - int(h*0.8)
    return _rand_pos(SW, SH, w, h)

# Data loader (COCO-Karpathy: images & captions as separate pools)
# Load fixed amount of COCO from Hugging Face and concatenate captions & images into lists.
def load_images_and_captions(
    split: str,
    images_needed: int,
    captions_needed: int,
    timeout: int,
    skip_broken: bool,
    max_fetch: int = 400
) -> Tuple[List[Image.Image], List[str]]:
    """
    Collect two *independent* pools:
      - images: list of PIL.Image (downloaded)
      - captions: list of strings (first caption)
    Returns when each pool reaches its target or dataset slice is exhausted.
    """
    hf_split = "train" if split == "train" else "val"
    # To avoid pulling the entire dataset, just slice the first `max_fetch` examples
    ds = load_dataset("yerevann/coco-karpathy", split=f"{hf_split}[:{max_fetch}]")

    def extract_caption(ex):
        if "sentences" in ex and ex["sentences"]:
            s0 = ex["sentences"][0]
            if isinstance(s0, dict):
                return s0.get("raw") or s0.get("caption") or s0.get("text") or ""
            elif isinstance(s0, str):
                return s0
        if "captions" in ex and ex["captions"]:
            c0 = ex["captions"][0]
            if isinstance(c0, dict):
                return c0.get("caption") or c0.get("raw") or c0.get("text") or ""
            elif isinstance(c0, str):
                return c0
        return ex.get("caption") or ex.get("text") or ""

    def extract_url(ex):
        return ex.get("coco_url") or ex.get("image_url") or ex.get("url")

    images: List[Image.Image] = []
    captions: List[str] = []
    seen_urls = set()

    # First pass: collect captions (no network)
    for ex in ds:
        if len(captions) >= captions_needed:
            break
        cap = (extract_caption(ex) or "").strip()
        if cap:
            captions.append(cap)
    # Second pass: collect images (network)
    for ex in ds:
        if len(images) >= images_needed:
            break
        url = extract_url(ex)
        if not url or url in seen_urls:
            continue
        try:
            resp = requests.get(url, timeout=timeout)
            resp.raise_for_status()
            img = Image.open(io.BytesIO(resp.content)).convert("RGB")
            images.append(img)
            seen_urls.add(url)
        except Exception:
            if not skip_broken:
                raise
            continue

    if len(images) < images_needed:
        print(f"[WARN] Only got {len(images)}/{images_needed} images from first {max_fetch} items.")
    if len(captions) < captions_needed:
        print(f"[WARN] Only got {len(captions)}/{captions_needed} captions from first {max_fetch} items.")

    return images, captions


# Generating mismatched & wrong-format slides
def build_mismatch_dirty_pptx(
    images: List[Image.Image],
    captions: List[str],
    out_path: str,
    overlap_prob: float = 0.75,
    breach_prob: float = 0.60,
    rotate_deg_max: float = 10.0,
    scale_range: Tuple[float, float] = (0.22, 0.60),
    font_candidates: Tuple[int, ...] = (26, 32, 36, 40),
    seed: int = 42
) -> str:
    """
    Create a deck where each slide has:
      - one random image (random size/pos, optional rotation)
      - one random caption (random size/pos/font, optional rotation)
    Image/Caption are mismatched by pairing shuffled pools (different orders).
    """
    random.seed(seed)

    # Shuffle separately; ensure mismatch by offset pairing
    n = min(len(images), len(captions))
    if n == 0:
        raise ValueError("No images or no captions to build slides.")
    img_idx = list(range(len(images)))
    cap_idx = list(range(len(captions)))
    random.shuffle(img_idx)
    random.shuffle(cap_idx)
    # If lengths equal, shift captions by a nonzero offset to avoid trivial pairing
    if len(cap_idx) == len(img_idx) and len(cap_idx) > 1:
        offset = random.randint(1, len(cap_idx)-1)
        cap_idx = cap_idx[offset:] + cap_idx[:offset]

    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = emu2px(prs.slide_width), emu2px(prs.slide_height)

    for i in range(n):
        img = images[img_idx[i % len(img_idx)]]
        cap = captions[cap_idx[i % len(cap_idx)]] or "No caption."

        slide = prs.slides.add_slide(blank)

        # random sizes
        iw, ih = _rand_size(SW, SH, scale_range)
        tw, th = _rand_size(SW, SH, scale_range)

        # random positions
        ix, iy = _rand_pos(SW, SH, iw, ih)
        tx, ty = _rand_pos(SW, SH, tw, th)

        # force overlap sometimes
        if random.random() < overlap_prob:
            tx, ty = _force_overlap(ix, iy, iw, ih, tw, th, SW, SH)

        # force margin breach sometimes
        if random.random() < breach_prob:
            side = random.choice(["left", "top", "right", "bottom"])
            if random.random() < 0.5:
                tx, ty = _force_margin_breach(tw, th, SW, SH, side)
            else:
                ix, iy = _force_margin_breach(iw, ih, SW, SH, side)

        # add image
        bio = io.BytesIO()
        img.convert("RGB").save(bio, format="JPEG", quality=90)
        bio.seek(0)
        pic = slide.shapes.add_picture(bio, px2emu(ix), px2emu(iy),
                                       width=px2emu(iw), height=px2emu(ih))

        # add caption textbox with big random font
        tb = slide.shapes.add_textbox(px2emu(tx), px2emu(ty), px2emu(tw), px2emu(th))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = cap
        r.font.size = Pt(random.choice(font_candidates))
        p.alignment = PP_ALIGN.LEFT

        # tiny rotation to look messier
        if rotate_deg_max > 0:
            try:
                pic.rotation = random.uniform(-rotate_deg_max, rotate_deg_max)
            except Exception:
                pass
            try:
                tb.rotation = random.uniform(-rotate_deg_max, rotate_deg_max)
            except Exception:
                pass

    prs.save(out_path)
    return out_path


# ---------- CLI ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--images", type=int, default=20, help="number of images to collect")
    ap.add_argument("--captions", type=int, default=20, help="number of captions to collect")
    ap.add_argument("--split", type=str, default="train", choices=["train", "validation"])
    ap.add_argument("--timeout", type=int, default=20, help="HTTP timeout (seconds)")
    ap.add_argument("--skip_broken", action="store_true", help="skip samples when an image fails to download")
    ap.add_argument("--max_fetch", type=int, default=400, help="max dataset examples to inspect")
    ap.add_argument("--out", type=str, default="coco_mismatch_dirty.pptx")

    # dirtiness knobs
    ap.add_argument("--overlap-prob", type=float, default=0.75)
    ap.add_argument("--breach-prob", type=float, default=0.60)
    ap.add_argument("--rotate-deg-max", type=float, default=10.0)
    ap.add_argument("--scale-range", type=parse_scale_range, default=(0.22, 0.60))
    ap.add_argument("--font-candidates", type=parse_font_candidates, default=(26, 32, 36, 40))
    ap.add_argument("--seed", type=int, default=42)

    args = ap.parse_args()

    imgs, caps = load_images_and_captions(
        split=args.split,
        images_needed=args.images,
        captions_needed=args.captions,
        timeout=args.timeout,
        skip_broken=args.skip_broken,
        max_fetch=args.max_fetch,
    )

    if not imgs or not caps:
        raise RuntimeError(f"Insufficient pools: images={len(imgs)}, captions={len(caps)}")

    out_path = Path(args.out).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    deck = build_mismatch_dirty_pptx(
        images=imgs,
        captions=caps,
        out_path=str(out_path),
        overlap_prob=args.overlap_prob,
        breach_prob=args.breach_prob,
        rotate_deg_max=args.rotate_deg_max,
        scale_range=args.scale_range,
        font_candidates=args.font_candidates,
        seed=args.seed,
    )
    print(f"[MISMATCH DIRTY] Saved: {deck}")

if __name__ == "__main__":
    main()
